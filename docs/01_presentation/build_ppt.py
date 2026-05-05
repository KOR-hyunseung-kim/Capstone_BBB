"""BBB 캡스톤 디자인 제안서 PowerPoint 생성 스크립트."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import pptx.oxml.ns as nsmap
from lxml import etree

# ── 색상 팔레트 ──────────────────────────────────────
BG        = RGBColor(0x05, 0x0a, 0x14)
SURFACE   = RGBColor(0x0d, 0x1b, 0x2e)
SURFACE2  = RGBColor(0x12, 0x20, 0x40)
ACCENT    = RGBColor(0x00, 0xc8, 0xff)
ACCENT2   = RGBColor(0x00, 0xff, 0x9d)
ACCENT3   = RGBColor(0x7b, 0x61, 0xff)
TEXT      = RGBColor(0xe8, 0xf0, 0xfe)
TEXT_DIM  = RGBColor(0x7a, 0x9a, 0xbf)
DANGER    = RGBColor(0xff, 0x4d, 0x6d)
WARN      = RGBColor(0xff, 0xd1, 0x66)
WHITE     = RGBColor(0xff, 0xff, 0xff)

W = Inches(13.33)   # 와이드 16:9
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

BLANK = prs.slide_layouts[6]   # 완전 빈 레이아웃

# ═══════════════════════════════════════════════════════════════
# 헬퍼 함수
# ═══════════════════════════════════════════════════════════════

def add_rect(slide, x, y, w, h, fill=None, line=None, line_w=Pt(1)):
    shape = slide.shapes.add_shape(1, x, y, w, h)   # MSO_SHAPE_TYPE.RECTANGLE = 1
    shape.line.width = line_w
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line:
        shape.line.color.rgb = line
    else:
        shape.line.fill.background()
    return shape


def add_text(slide, text, x, y, w, h,
             size=Pt(14), bold=False, color=TEXT, align=PP_ALIGN.LEFT,
             wrap=True, italic=False):
    txb = slide.shapes.add_textbox(x, y, w, h)
    tf  = txb.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = size
    run.font.bold  = bold
    run.font.color.rgb = color
    run.font.italic = italic
    return txb


def bg(slide):
    """전체 배경을 진한 네이비로."""
    add_rect(slide, 0, 0, W, H, fill=BG)


def slide_label(slide, text):
    add_text(slide, text,
             Inches(0.6), Inches(0.3), Inches(8), Inches(0.35),
             size=Pt(10), bold=True, color=ACCENT, align=PP_ALIGN.LEFT)


def section_title(slide, text, y=Inches(0.75)):
    add_text(slide, text,
             Inches(0.6), y, Inches(12), Inches(0.7),
             size=Pt(36), bold=True, color=TEXT, align=PP_ALIGN.LEFT)


def accent_bar(slide, y=Inches(1.55)):
    """제목 아래 가로선."""
    add_rect(slide, Inches(0.6), y, Inches(12.1), Inches(0.03), fill=ACCENT2)


def card(slide, x, y, w, h, border=SURFACE2):
    add_rect(slide, x, y, w, h, fill=SURFACE, line=border, line_w=Pt(1.2))


def progress_bar(slide, current, total):
    pct = current / total
    add_rect(slide, 0, H - Inches(0.06), W, Inches(0.06), fill=SURFACE2)
    add_rect(slide, 0, H - Inches(0.06), Emu(int(W * pct)), Inches(0.06), fill=ACCENT)
    add_text(slide, f"{current} / {total}",
             W - Inches(1.2), H - Inches(0.35),
             Inches(1.1), Inches(0.3),
             size=Pt(9), color=TEXT_DIM, align=PP_ALIGN.RIGHT)


# ═══════════════════════════════════════════════════════════════
# SLIDE 1 — 타이틀
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)

# 대형 BBB 로고
add_text(s, "BBB",
         Inches(0.6), Inches(1.2), Inches(12.1), Inches(2.4),
         size=Pt(120), bold=True, color=ACCENT, align=PP_ALIGN.CENTER)

# Bio Body Band
add_text(s, "Bio Body Band",
         Inches(0.6), Inches(3.3), Inches(12.1), Inches(0.55),
         size=Pt(22), bold=False, color=TEXT_DIM, align=PP_ALIGN.CENTER)

# 설명
add_text(s,
         "EMG + IMU 센서로 작업자의 근피로를 모니터링하고,\n"
         "오염 환경에서도 블루투스 핸즈프리 제어를 실현하는 스마트 웨어러블 아대",
         Inches(2.0), Inches(4.0), Inches(9.3), Inches(0.9),
         size=Pt(14), color=TEXT_DIM, align=PP_ALIGN.CENTER)

# 태그 배지들
tags = [("ESP32-S3", ACCENT), ("MicroPython", ACCENT2),
        ("BLE HID", ACCENT3), ("EMG · IMU", DANGER)]
for i, (t, c) in enumerate(tags):
    bx = Inches(3.0 + i * 1.9)
    add_rect(s, bx, Inches(5.05), Inches(1.7), Inches(0.38),
             fill=SURFACE2, line=c, line_w=Pt(1.2))
    add_text(s, t, bx, Inches(5.05), Inches(1.7), Inches(0.38),
             size=Pt(11), bold=True, color=c, align=PP_ALIGN.CENTER)

# 슬로건
add_text(s, '"Protecting the Body, Bridging the Device."',
         Inches(0.6), Inches(5.7), Inches(12.1), Inches(0.4),
         size=Pt(11), italic=True, color=TEXT_DIM, align=PP_ALIGN.CENTER)

# Capstone Design 2026
add_text(s, "Capstone Design 2026",
         Inches(0.6), Inches(6.2), Inches(12.1), Inches(0.4),
         size=Pt(10), color=TEXT_DIM, align=PP_ALIGN.CENTER)

progress_bar(s, 1, 9)


# ═══════════════════════════════════════════════════════════════
# SLIDE 2 — 왜 BBB가 필요한가?
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
slide_label(s, "Problem Statement")
section_title(s, "왜 BBB가 필요한가?")
accent_bar(s)

add_text(s, "기존 안전 솔루션이 해결하지 못한 3가지 문제",
         Inches(0.6), Inches(1.65), Inches(12), Inches(0.4),
         size=Pt(13), color=TEXT_DIM)

problems = [
    ("⚠️  근피로 데이터 부재",
     "기존 안전 장구는 외부 충격 보호에만 집중.\n"
     "근피로 누적 → 물건 낙하·중대 사고 위험.\n"
     "이를 실시간으로 감지·경보하는 웨어러블 장치가 없음."),
    ("🧤  현장 오염 문제",
     "기름·화학약품이 묻은 손으로 노트북·태블릿 조작\n"
     "→ 정밀 기기 오염·고장.\n"
     "비접촉 입력 수단이 절실."),
    ("🤚  비접촉식 입력장치 부재",
     "신체 동작만으로 기기를 직접 만지지 않고 제어할 수 있는\n"
     "산업 현장용 비접촉식 입력장치가 없음.\n"
     "손을 쓰지 못하는 상황에서도 즉시 조작 가능해야 함."),
]

for i, (title, body) in enumerate(problems):
    cx = Inches(0.55 + i * 4.25)
    card(s, cx, Inches(2.15), Inches(4.1), Inches(4.8), border=SURFACE2)
    add_text(s, title, cx + Inches(0.2), Inches(2.3), Inches(3.7), Inches(0.5),
             size=Pt(14), bold=True, color=TEXT)
    add_text(s, body, cx + Inches(0.2), Inches(2.85), Inches(3.7), Inches(3.8),
             size=Pt(12), color=TEXT_DIM)

progress_bar(s, 2, 9)


# ═══════════════════════════════════════════════════════════════
# SLIDE 3 — 프로젝트 목표
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
slide_label(s, "Project Goals")
section_title(s, "프로젝트 목표")
accent_bar(s)

# Safety Mode 카드
card(s, Inches(0.55), Inches(1.7), Inches(6.0), Inches(5.3), border=DANGER)
add_text(s, "🛡️  Safety Mode", Inches(0.75), Inches(1.9), Inches(5.6), Inches(0.5),
         size=Pt(16), bold=True, color=DANGER)
add_text(s, "근피로 모니터링", Inches(0.75), Inches(2.45), Inches(5.6), Inches(0.4),
         size=Pt(20), bold=True, color=TEXT)
safety_items = [
    "EMG Median Frequency 실시간 분석",
    "피로 임계치 80% / 95% 이중 경보",
    "진동 + LED 촉각·시각 피드백",
    "근골격계 질환 사전 예방",
]
for j, item in enumerate(safety_items):
    add_text(s, f"→  {item}", Inches(0.75), Inches(3.0 + j * 0.62), Inches(5.5), Inches(0.5),
             size=Pt(13), color=TEXT_DIM)

# Control Mode 카드
card(s, Inches(6.78), Inches(1.7), Inches(6.0), Inches(5.3), border=ACCENT)
add_text(s, "🎮  Control Mode", Inches(6.98), Inches(1.9), Inches(5.6), Inches(0.5),
         size=Pt(16), bold=True, color=ACCENT)
add_text(s, "핸즈프리 BLE 마우스", Inches(6.98), Inches(2.45), Inches(5.6), Inches(0.4),
         size=Pt(20), bold=True, color=TEXT)
control_items = [
    "IMU 팔 기울기 → 커서 이동",
    "EMG Spike (주먹 쥐기) → 클릭",
    "BLE HID — 드라이버 없이 즉시 연결",
    "Windows · Android · iOS 호환",
]
for j, item in enumerate(control_items):
    add_text(s, f"→  {item}", Inches(6.98), Inches(3.0 + j * 0.62), Inches(5.5), Inches(0.5),
             size=Pt(13), color=TEXT_DIM)

progress_bar(s, 3, 9)


# ═══════════════════════════════════════════════════════════════
# SLIDE 4 — 하드웨어 구성
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
slide_label(s, "Hardware Architecture")
section_title(s, "하드웨어 구성")
accent_bar(s)

add_text(s, "설계 원칙: 경량화 · 소형화 최우선  —  불필요한 부품 제거, 전류 예산 내 설계",
         Inches(0.6), Inches(1.65), Inches(12), Inches(0.38),
         size=Pt(12), color=ACCENT2)

hw_items = [
    ("🧠", "MCU",       "XIAO ESP32-S3",       "BLE HID · WiFi · USB-C 충전 내장",  False),
    ("💪", "EMG",       "MyoWare 2.0",         "근육 수축 강도 · 피로도 수집",        False),
    ("📐", "IMU",       "MPU6050",             "팔 기울기 X/Y · 6축 관성 감지",       False),
    ("🔋", "Battery",   "Li-Po 3.7V 300mAh",  "초슬림 · 2시간+ 연속 동작",           False),
    ("📳", "Vibration", "Coin Motor",          "근위험 도달 시 촉각 경보",            True),
    ("🖥️", "Display",   "OLED 0.96\"",         "모드 · 배터리 시각화",               True),
    ("🩹", "Electrode", "Disposable Patch",    "피부 접촉용 교체형 전극",             False),
    ("📦", "Housing",   "3D Print FDM",        "PLA/PETG · 초소형 경량 설계",        True),
]

cols, rows = 4, 2
cw, ch = Inches(3.05), Inches(2.1)
for idx, (icon, name, model, role, optional) in enumerate(hw_items):
    col, row = idx % cols, idx // cols
    cx = Inches(0.35 + col * 3.17)
    cy = Inches(2.15 + row * 2.25)
    border = WARN if optional else SURFACE2
    card(s, cx, cy, cw, ch, border=border)
    add_text(s, icon, cx + Inches(0.15), cy + Inches(0.12), Inches(0.5), Inches(0.5),
             size=Pt(20))
    lbl = name + ("  [Option]" if optional else "")
    add_text(s, lbl, cx + Inches(0.65), cy + Inches(0.14), Inches(2.2), Inches(0.38),
             size=Pt(13), bold=True, color=WARN if optional else TEXT)
    add_text(s, model, cx + Inches(0.15), cy + Inches(0.62), Inches(2.8), Inches(0.32),
             size=Pt(10), bold=True, color=ACCENT2)
    add_text(s, role, cx + Inches(0.15), cy + Inches(1.0), Inches(2.8), Inches(0.85),
             size=Pt(10), color=TEXT_DIM)

progress_bar(s, 4, 9)


# ═══════════════════════════════════════════════════════════════
# SLIDE 5 — 동작 흐름
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
slide_label(s, "System Flow")
section_title(s, "아대 동작 흐름")
accent_bar(s)

# Safety Flow
add_text(s, "Safety Mode  —  근피로 모니터링 파이프라인",
         Inches(0.6), Inches(1.7), Inches(12), Inches(0.38),
         size=Pt(12), bold=True, color=DANGER)

safety_nodes = ["EMG 샘플링\n1kHz ADC", "Band-pass\n20~500Hz",
                "FFT 분석\nMedian Freq.", "80% 경고\n진동 1회", "95% 긴급\n강진동+LED"]
node_colors  = [ACCENT, TEXT, TEXT, WARN, DANGER]
for i, (node, nc) in enumerate(zip(safety_nodes, node_colors)):
    nx = Inches(0.55 + i * 2.5)
    card(s, nx, Inches(2.15), Inches(2.1), Inches(0.85), border=nc)
    add_text(s, node, nx, Inches(2.15), Inches(2.1), Inches(0.85),
             size=Pt(10), bold=True, color=nc, align=PP_ALIGN.CENTER)
    if i < 4:
        add_text(s, "→", Inches(2.55 + i * 2.5), Inches(2.35), Inches(0.45), Inches(0.4),
                 size=Pt(18), color=TEXT_DIM, align=PP_ALIGN.CENTER)

# Control Flow
add_text(s, "Control Mode  —  BLE HID 마우스 파이프라인",
         Inches(0.6), Inches(3.3), Inches(12), Inches(0.38),
         size=Pt(12), bold=True, color=ACCENT)

ctrl_nodes  = ["IMU 읽기\n100Hz I2C", "Comp. Filter\npitch / roll",
               "커서 이동\nBLE HID ΔX/Y", "EMG Spike\n주먹 쥐기", "BLE 클릭\nHID Button"]
ctrl_colors = [ACCENT, TEXT, ACCENT2, TEXT, ACCENT2]
for i, (node, nc) in enumerate(zip(ctrl_nodes, ctrl_colors)):
    nx = Inches(0.55 + i * 2.5)
    card(s, nx, Inches(3.75), Inches(2.1), Inches(0.85), border=nc)
    add_text(s, node, nx, Inches(3.75), Inches(2.1), Inches(0.85),
             size=Pt(10), bold=True, color=nc, align=PP_ALIGN.CENTER)
    if i < 4:
        add_text(s, "→", Inches(2.55 + i * 2.5), Inches(3.95), Inches(0.45), Inches(0.4),
                 size=Pt(18), color=TEXT_DIM, align=PP_ALIGN.CENTER)

# 레이턴시 라벨
add_rect(s, Inches(4.5), Inches(5.1), Inches(4.3), Inches(0.4), fill=SURFACE2)
add_text(s, "전체 파이프라인 레이턴시 목표: < 50ms",
         Inches(4.5), Inches(5.1), Inches(4.3), Inches(0.4),
         size=Pt(11), color=ACCENT2, align=PP_ALIGN.CENTER)

progress_bar(s, 5, 9)


# ═══════════════════════════════════════════════════════════════
# SLIDE 6 — 시나리오 A: Safety Mode
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
slide_label(s, "Demo Scenario A")

add_text(s, "Safety Mode", Inches(0.6), Inches(0.7), Inches(6), Inches(0.6),
         size=Pt(36), bold=True, color=TEXT)
add_rect(s, Inches(0.6), Inches(1.55), Inches(12.1), Inches(0.03), fill=DANGER)

add_text(s, "근골격계 질환 예방  |  물류 운반 작업자 시나리오",
         Inches(0.6), Inches(1.65), Inches(12), Inches(0.38),
         size=Pt(12), color=TEXT_DIM)

steps_a = [
    ("1", "착용 & 시작",
     "아대 착용 후 전원 ON → BLE로 모니터링 앱에 자동 연결"),
    ("2", "피로도 누적 감지",
     "반복 작업 → EMG Median Frequency가 점진적으로 저주파 방향 이동"),
    ("3", "1단계 경보 — 80%",
     "진동 1회 + LED 노란색 점등 → \"휴식 권장\" 신호 발생"),
    ("4", "결과",
     "작업자가 알람으로 휴식 타이밍 인지 → 누적 과부하 차단 → 근골격계 질환 예방"),
]
for i, (num, title, body) in enumerate(steps_a):
    cy = Inches(2.2 + i * 1.15)
    # 번호 원
    add_rect(s, Inches(0.55), cy + Inches(0.05), Inches(0.45), Inches(0.45),
             fill=DANGER, line=DANGER)
    add_text(s, num, Inches(0.55), cy + Inches(0.05), Inches(0.45), Inches(0.45),
             size=Pt(12), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    card(s, Inches(1.15), cy, Inches(11.5), Inches(1.0), border=SURFACE2)
    add_text(s, title, Inches(1.35), cy + Inches(0.1), Inches(4), Inches(0.38),
             size=Pt(14), bold=True, color=TEXT)
    add_text(s, body, Inches(1.35), cy + Inches(0.48), Inches(11.0), Inches(0.45),
             size=Pt(12), color=TEXT_DIM)

progress_bar(s, 6, 9)


# ═══════════════════════════════════════════════════════════════
# SLIDE 7 — 시나리오 B: Control Mode
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
slide_label(s, "Demo Scenario B")

add_text(s, "Control Mode", Inches(0.6), Inches(0.7), Inches(7), Inches(0.6),
         size=Pt(36), bold=True, color=TEXT)
add_rect(s, Inches(0.6), Inches(1.55), Inches(12.1), Inches(0.03), fill=ACCENT)

add_text(s, "오염 환경 핸즈프리 제어  |  유해물질 접촉 정비사 시나리오",
         Inches(0.6), Inches(1.65), Inches(12), Inches(0.38),
         size=Pt(12), color=TEXT_DIM)

steps_b = [
    ("1", "상황",
     "장갑에 오일이 묻은 정비사 → 매뉴얼 확인을 위해 노트북 앞에 서 있음"),
    ("2", "커서 이동 (IMU)",
     "팔을 상·하·좌·우로 기울이면 pitch/roll이 BLE HID ΔX/Y로 변환 → 커서 이동"),
    ("3", "클릭 (EMG Spike)",
     "주먹을 순간적으로 꽉 쥐면 EMG 스파이크 감지 → 왼쪽 클릭 → 도면 열기·확대"),
    ("4", "결과",
     "기기를 한 번도 건드리지 않고 정보 확인 완료 → 오염 전파 Zero, 장비 고장 방지"),
]
for i, (num, title, body) in enumerate(steps_b):
    cy = Inches(2.2 + i * 1.15)
    add_rect(s, Inches(0.55), cy + Inches(0.05), Inches(0.45), Inches(0.45),
             fill=ACCENT, line=ACCENT)
    add_text(s, num, Inches(0.55), cy + Inches(0.05), Inches(0.45), Inches(0.45),
             size=Pt(12), bold=True, color=BG, align=PP_ALIGN.CENTER)
    card(s, Inches(1.15), cy, Inches(11.5), Inches(1.0), border=SURFACE2)
    add_text(s, title, Inches(1.35), cy + Inches(0.1), Inches(4), Inches(0.38),
             size=Pt(14), bold=True, color=TEXT)
    add_text(s, body, Inches(1.35), cy + Inches(0.48), Inches(11.0), Inches(0.45),
             size=Pt(12), color=TEXT_DIM)

progress_bar(s, 7, 9)


# ═══════════════════════════════════════════════════════════════
# SLIDE 8 — 개발 일정
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
slide_label(s, "Development Schedule")
section_title(s, "12주 개발 일정")
accent_bar(s)

milestones = [
    ("W1–W2",  "환경 구축",        "부품 수급 · ESP32-S3 BLE 기본 통신 환경 구축",               ACCENT),
    ("W3–W5",  "신호처리 알고리즘", "EMG Band-pass 필터 · FFT Median Frequency · IMU Comp. Filter", ACCENT3),
    ("W6–W8",  "BLE HID + 하우징", "Bluetooth HID 마우스 프로토콜 최적화 · 3D 프린팅 경량 하우징",  ACCENT2),
    ("W9–W11", "통합 테스트 & 보정", "시나리오별 데이터 보정 · 사용자 피드백 반영 · QA",            WARN),
    ("W12",    "최종 발표 & 시연",  "2026년 6월 12일 — Safety & Control 두 시나리오 라이브 데모",   DANGER),
]

for i, (week, title, body, color) in enumerate(milestones):
    cy = Inches(1.75 + i * 1.08)
    # 날짜 칩
    add_rect(s, Inches(0.55), cy, Inches(1.1), Inches(0.42), fill=SURFACE2)
    add_text(s, week, Inches(0.55), cy, Inches(1.1), Inches(0.42),
             size=Pt(10), bold=True, color=color, align=PP_ALIGN.CENTER)
    # 연결선
    add_rect(s, Inches(1.65), cy + Inches(0.18), Inches(0.3), Inches(0.06), fill=color)
    # 카드
    card(s, Inches(1.95), cy, Inches(10.8), Inches(0.9), border=color)
    add_text(s, title, Inches(2.15), cy + Inches(0.06), Inches(4), Inches(0.35),
             size=Pt(13), bold=True, color=TEXT)
    add_text(s, body, Inches(2.15), cy + Inches(0.44), Inches(10.3), Inches(0.38),
             size=Pt(11), color=TEXT_DIM)

progress_bar(s, 8, 9)


# ═══════════════════════════════════════════════════════════════
# SLIDE 9 — 기대 효과 & 마무리
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
slide_label(s, "Impact & Closing")
section_title(s, "기대 효과")
accent_bar(s)

impacts = [
    ("🏭", "산업 안전 향상",     "근피로 기반 휴식 타이밍 제공\n→ 근골격계 질환 발생률 감소"),
    ("🔬", "오염 환경 생산성",   "기기 비접촉 제어 → 장비 오염·고장 Zero\n→ 정비·연구 현장 적용"),
    ("📡", "표준 BLE HID",      "드라이버 없이 Win/Android/iOS\n즉시 연결 — 범용성 확보"),
    ("⚡", "경량 웨어러블",      "초소형 폼팩터 + Li-Po 2h+\n→ 장시간 착용 현실적"),
    ("🧬", "바이오 신호처리",    "EMG FFT + IMU 융합 알고리즘\n→ 학술적 기여 및 확장 플랫폼"),
    ("🏆", "최종 목표",          "2026.06.12 라이브 데모\nSafety & Control 완전 동작"),
]

cols_i = 3
for idx, (icon, title, body) in enumerate(impacts):
    col, row = idx % cols_i, idx // cols_i
    cx = Inches(0.5 + col * 4.27)
    cy = Inches(1.85 + row * 2.15)
    border = ACCENT if idx == 5 else SURFACE2
    card(s, cx, cy, Inches(4.0), Inches(1.95), border=border)
    add_text(s, icon, cx + Inches(0.18), cy + Inches(0.15), Inches(0.55), Inches(0.55),
             size=Pt(22))
    add_text(s, title, cx + Inches(0.75), cy + Inches(0.18), Inches(3.1), Inches(0.4),
             size=Pt(13), bold=True, color=ACCENT if idx == 5 else TEXT)
    add_text(s, body, cx + Inches(0.18), cy + Inches(0.7), Inches(3.65), Inches(1.1),
             size=Pt(11), color=TEXT_DIM)

# 슬로건
add_text(s, '"Protecting the Body, Bridging the Device."',
         Inches(0.6), Inches(6.4), Inches(12.1), Inches(0.45),
         size=Pt(14), italic=True, color=TEXT_DIM, align=PP_ALIGN.CENTER)

progress_bar(s, 9, 9)


# ─── 저장 ───────────────────────────────────────────────────────
OUT = "C:/4_1/BBB/docs/BBB_proposal.pptx"
prs.save(OUT)
print(f"Done: {OUT}")
